from typing import Type
from crewai_tools import BaseTool
from pydantic import BaseModel, Field
from pptx import Presentation


class PowerPointInput(BaseModel):
    """Input model for the PowerPoint Tool."""
    data: dict = Field(..., description="Dictionary of information about the destination to be passed to the tool.")
    output_file: str = Field(default="Presentation.pptx", description="Name of the output PowerPoint file.")


class PowerPointTool(BaseTool):
    """Tool to create a PowerPoint presentation."""
    name: str = "PowerPoint Tool"
    description: str = "Create a PowerPoint presentation from a dictionary of data."
    args_schema: Type[BaseModel] = PowerPointInput

    def _run(self, data: dict, output_file: str = "Presentation.pptx") -> str:
        """Run the tool with provided data."""

        prs = Presentation()
        
        # Title Slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        title.text = "Travel Presentation"
        subtitle.text = "Your custom travel itinerary"

        # Add slides for each data category
        for section, content in data.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section
            slide.placeholders[1].text = content

        # Save the presentation
        prs.save(output_file)
        return output_file